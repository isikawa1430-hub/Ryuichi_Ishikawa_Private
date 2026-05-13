from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Color palette
COLOR_NAVY = RGBColor(0x1A, 0x37, 0x6C)      # Dark navy blue (primary)
COLOR_ACCENT = RGBColor(0x00, 0x7B, 0xC2)     # Accent blue
COLOR_LIGHT_BLUE = RGBColor(0xE8, 0xF4, 0xFD) # Light blue background
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
COLOR_MID_GRAY = RGBColor(0x66, 0x66, 0x66)
COLOR_LIGHT_GRAY = RGBColor(0xF0, 0xF4, 0xF8)
COLOR_ORANGE = RGBColor(0xE8, 0x7A, 0x1E)     # Accent orange
COLOR_GREEN = RGBColor(0x2E, 0x86, 0x48)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = 6


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=COLOR_DARK_GRAY,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Meiryo UI"
    return txBox


def add_paragraph(tf, text, font_size=14, bold=False, color=COLOR_DARK_GRAY,
                  align=PP_ALIGN.LEFT, space_before=Pt(4), indent_level=0, italic=False):
    from pptx.util import Pt as Pt2
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    p.level = indent_level
    run = p.add_run()
    run.text = text
    run.font.size = Pt2(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Meiryo UI"
    return p


def set_first_paragraph(tf, text, font_size=14, bold=False, color=COLOR_DARK_GRAY,
                         align=PP_ALIGN.LEFT, italic=False):
    p = tf.paragraphs[0]
    p.alignment = align
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Meiryo UI"


def add_slide_header(slide, title, subtitle=None):
    """Add top header bar with title"""
    # Header background bar
    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.1), fill_color=COLOR_NAVY)
    # Accent line
    add_rect(slide, Inches(0), Inches(1.1), Inches(13.33), Inches(0.05), fill_color=COLOR_ACCENT)

    # Title text
    tb = add_text_box(slide, title, Inches(0.4), Inches(0.1), Inches(10), Inches(0.9),
                      font_size=24, bold=True, color=COLOR_WHITE, align=PP_ALIGN.LEFT)

    if subtitle:
        add_text_box(slide, subtitle, Inches(0.4), Inches(0.75), Inches(12), Inches(0.4),
                     font_size=12, bold=False, color=RGBColor(0xB0, 0xC8, 0xE8), align=PP_ALIGN.LEFT)

    # Slide number / footer
    add_text_box(slide, "株式会社梓総合研究所  |  AIR-Plate 事業提携ご提案",
                 Inches(0.3), Inches(7.15), Inches(10), Inches(0.3),
                 font_size=8, color=COLOR_MID_GRAY)


def add_section_label(slide, text, left, top, width=Inches(3)):
    """Add colored section label"""
    add_rect(slide, left, top, width, Inches(0.32), fill_color=COLOR_ACCENT)
    add_text_box(slide, text, left + Inches(0.1), top, width - Inches(0.1), Inches(0.32),
                 font_size=11, bold=True, color=COLOR_WHITE)


def add_bullet_box(slide, items, left, top, width, height,
                   font_size=13, header=None, header_color=COLOR_NAVY,
                   bg_color=COLOR_LIGHT_GRAY, bullet_char="●"):
    """Add a bullet point box with optional header"""
    add_rect(slide, left, top, width, height, fill_color=bg_color)

    y_offset = top + Inches(0.12)
    if header:
        add_text_box(slide, header, left + Inches(0.15), y_offset, width - Inches(0.2), Inches(0.35),
                     font_size=12, bold=True, color=header_color)
        y_offset += Inches(0.35)

    item_height = (height - (y_offset - top) - Inches(0.1)) / max(len(items), 1)
    for item in items:
        add_text_box(slide, f"{bullet_char}  {item}", left + Inches(0.15), y_offset,
                     width - Inches(0.3), max(item_height, Inches(0.35)),
                     font_size=font_size, color=COLOR_DARK_GRAY)
        y_offset += item_height


# ─────────────────────────────────────────
# SLIDE 1: Cover
# ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])

# Full background
add_rect(slide, 0, 0, Inches(13.33), Inches(7.5), fill_color=COLOR_NAVY)

# Right accent panel
add_rect(slide, Inches(9.5), 0, Inches(3.83), Inches(7.5), fill_color=COLOR_ACCENT)

# Diagonal accent
add_rect(slide, Inches(8.8), 0, Inches(0.8), Inches(7.5), fill_color=RGBColor(0x00, 0x5E, 0x9A))

# Top line
add_rect(slide, 0, Inches(0.5), Inches(9.0), Inches(0.04), fill_color=COLOR_ORANGE)

# Main title
add_text_box(slide, "AIR-Plate", Inches(0.5), Inches(1.2), Inches(8.5), Inches(1.4),
             font_size=54, bold=True, color=COLOR_WHITE, align=PP_ALIGN.LEFT)

add_text_box(slide, "事業提携のご提案", Inches(0.5), Inches(2.5), Inches(8.5), Inches(1.0),
             font_size=32, bold=True, color=COLOR_WHITE, align=PP_ALIGN.LEFT)

add_text_box(slide, "― 販売代理店契約に向けた協業フレームワーク ―",
             Inches(0.5), Inches(3.4), Inches(8.5), Inches(0.6),
             font_size=16, bold=False, color=RGBColor(0xB0, 0xC8, 0xE8), align=PP_ALIGN.LEFT)

# Divider line
add_rect(slide, Inches(0.5), Inches(4.1), Inches(6.0), Inches(0.03), fill_color=COLOR_ORANGE)

# Company & date
add_text_box(slide, "株式会社梓総合研究所", Inches(0.5), Inches(4.25), Inches(8), Inches(0.5),
             font_size=18, bold=True, color=COLOR_WHITE)
add_text_box(slide, "2026年5月", Inches(0.5), Inches(4.75), Inches(8), Inches(0.4),
             font_size=14, color=RGBColor(0xB0, 0xC8, 0xE8))

# Right panel content
add_text_box(slide, "CCT様向け\n打合せ資料", Inches(9.7), Inches(2.5), Inches(3.2), Inches(1.5),
             font_size=20, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, "Confidential", Inches(9.7), Inches(6.8), Inches(3.2), Inches(0.4),
             font_size=10, color=RGBColor(0xD0, 0xE8, 0xFF), align=PP_ALIGN.CENTER, italic=True)


# ─────────────────────────────────────────
# SLIDE 2: Agenda
# ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
add_rect(slide, 0, 0, Inches(13.33), Inches(7.5), fill_color=RGBColor(0xF7, 0xF9, 0xFC))
add_slide_header(slide, "本日のアジェンダ")

agenda_items = [
    ("1", "AIR-Plate事業の現状と成長機会"),
    ("2", "協業スキームのご提案（役割分担）"),
    ("3", "販売・収益モデル"),
    ("4", "案件運営の進め方"),
    ("5", "知財・秘密保持の考え方"),
    ("6", "契約の基本条件"),
    ("7", "本日ご確認いただきたいポイント"),
]

col1_items = agenda_items[:4]
col2_items = agenda_items[4:]

for i, (num, text) in enumerate(agenda_items):
    y = Inches(1.5) + i * Inches(0.72)
    # Number badge
    badge = add_rect(slide, Inches(0.5), y, Inches(0.45), Inches(0.45), fill_color=COLOR_NAVY)
    add_text_box(slide, num, Inches(0.5), y, Inches(0.45), Inches(0.45),
                 font_size=16, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    # Item text
    add_rect(slide, Inches(1.1), y, Inches(11.5), Inches(0.45),
             fill_color=COLOR_WHITE if i % 2 == 0 else COLOR_LIGHT_GRAY)
    add_text_box(slide, text, Inches(1.25), y + Inches(0.04), Inches(11.0), Inches(0.42),
                 font_size=16, color=COLOR_DARK_GRAY)
    # Accent line on left of item
    add_rect(slide, Inches(1.1), y, Inches(0.04), Inches(0.45), fill_color=COLOR_ACCENT)


# ─────────────────────────────────────────
# SLIDE 3: AIR-Plate Overview
# ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
add_rect(slide, 0, 0, Inches(13.33), Inches(7.5), fill_color=RGBColor(0xF7, 0xF9, 0xFC))
add_slide_header(slide, "AIR-Plateの概要と市場機会",
                 "建築DX・FM領域に特化した維持管理BIMプラットフォーム")

# Left: What is AIR-Plate
add_rect(slide, Inches(0.3), Inches(1.4), Inches(5.8), Inches(5.7), fill_color=COLOR_WHITE)
add_rect(slide, Inches(0.3), Inches(1.4), Inches(5.8), Inches(0.45), fill_color=COLOR_NAVY)
add_text_box(slide, "AIR-Plateとは", Inches(0.4), Inches(1.4), Inches(5.6), Inches(0.45),
             font_size=14, bold=True, color=COLOR_WHITE)

features = [
    "建築設計・FM領域に特化した\nNotionベースの維持管理BIMプラットフォーム",
    "3DスキャンデータとBIMデータ、\n業務テンプレートを統合",
    "建物のライフサイクル全体を\nデータで一元支援",
]
for i, feat in enumerate(features):
    y = Inches(2.0) + i * Inches(1.3)
    add_rect(slide, Inches(0.45), y, Inches(0.35), Inches(0.35), fill_color=COLOR_ACCENT)
    add_text_box(slide, str(i+1), Inches(0.45), y, Inches(0.35), Inches(0.35),
                 font_size=12, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, feat, Inches(0.9), y - Inches(0.05), Inches(5.0), Inches(0.6),
                 font_size=13, color=COLOR_DARK_GRAY)

# Middle: Market Growth
add_rect(slide, Inches(6.4), Inches(1.4), Inches(3.3), Inches(2.7), fill_color=COLOR_WHITE)
add_rect(slide, Inches(6.4), Inches(1.4), Inches(3.3), Inches(0.45), fill_color=COLOR_GREEN)
add_text_box(slide, "市場の成長性", Inches(6.5), Inches(1.4), Inches(3.1), Inches(0.45),
             font_size=14, bold=True, color=COLOR_WHITE)
growth_items = [
    "建築DX・BIM活用ニーズの拡大",
    "FM領域でのデータ\n一元管理ニーズの高まり",
    "自治体・大手デベロッパーからの\n引き合い増加",
]
for i, item in enumerate(growth_items):
    y = Inches(2.0) + i * Inches(0.7)
    add_rect(slide, Inches(6.55), y + Inches(0.08), Inches(0.08), Inches(0.08), fill_color=COLOR_GREEN)
    add_text_box(slide, item, Inches(6.75), y, Inches(2.8), Inches(0.65), font_size=12, color=COLOR_DARK_GRAY)

# Right: Challenge
add_rect(slide, Inches(9.9), Inches(1.4), Inches(3.1), Inches(2.7), fill_color=COLOR_WHITE)
add_rect(slide, Inches(9.9), Inches(1.4), Inches(3.1), Inches(0.45), fill_color=COLOR_ORANGE)
add_text_box(slide, "課題 ― パートナーの必要性", Inches(10.0), Inches(1.4), Inches(2.9), Inches(0.45),
             font_size=13, bold=True, color=COLOR_WHITE)
challenge_items = [
    "梓総研単独では\n営業・導入リソースに限界",
    "ITコンサルティング力と\n営業基盤を持つパートナーが不可欠",
]
for i, item in enumerate(challenge_items):
    y = Inches(2.0) + i * Inches(0.9)
    add_rect(slide, Inches(10.05), y + Inches(0.08), Inches(0.08), Inches(0.08), fill_color=COLOR_ORANGE)
    add_text_box(slide, item, Inches(10.25), y, Inches(2.6), Inches(0.8), font_size=12, color=COLOR_DARK_GRAY)

# Bottom summary box
add_rect(slide, Inches(0.3), Inches(5.4), Inches(12.7), Inches(1.8), fill_color=COLOR_LIGHT_BLUE)
add_rect(slide, Inches(0.3), Inches(5.4), Inches(0.07), Inches(1.8), fill_color=COLOR_ACCENT)
add_text_box(slide, "AIR-Plateは建物の「生まれてから終わりまで」をデジタルで繋ぐ、建築業界の新インフラです。\nITコンサルティング力と営業基盤を持つCCT様との協業で、市場浸透を加速させます。",
             Inches(0.55), Inches(5.5), Inches(12.3), Inches(1.6),
             font_size=13, color=COLOR_NAVY, italic=True)


# ─────────────────────────────────────────
# SLIDE 4: Why CCT Partnership
# ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
add_rect(slide, 0, 0, Inches(13.33), Inches(7.5), fill_color=RGBColor(0xF7, 0xF9, 0xFC))
add_slide_header(slide, "なぜCCT様とのパートナーシップか",
                 "CCT様の強み × AIR-Plateの親和性")

# CCT Strengths
add_rect(slide, Inches(0.3), Inches(1.4), Inches(7.8), Inches(3.5), fill_color=COLOR_WHITE)
add_rect(slide, Inches(0.3), Inches(1.4), Inches(7.8), Inches(0.45), fill_color=COLOR_ACCENT)
add_text_box(slide, "CCT様の強み × AIR-Plateの親和性", Inches(0.4), Inches(1.4), Inches(7.6), Inches(0.45),
             font_size=14, bold=True, color=COLOR_WHITE)

strengths = [
    ("ITコンサルティング力", "エンドユーザーの業務課題を構造的に理解し、最適な導入を実現できる"),
    ("Notion代理店としての基盤", "AIR-Plate（Notionベース）の販売と一体化した提案が可能"),
    ("営業・サポート体制", "大規模案件の運用・保守を担える組織力"),
]
for i, (title, desc) in enumerate(strengths):
    y = Inches(2.0) + i * Inches(0.9)
    add_rect(slide, Inches(0.45), y, Inches(2.0), Inches(0.7), fill_color=COLOR_LIGHT_BLUE)
    add_text_box(slide, title, Inches(0.5), y + Inches(0.05), Inches(1.9), Inches(0.6),
                 font_size=11, bold=True, color=COLOR_NAVY, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(2.5), y + Inches(0.1), Inches(5.4), Inches(0.5), fill_color=RGBColor(0xF0, 0xF4, 0xF8))
    add_text_box(slide, desc, Inches(2.6), y + Inches(0.05), Inches(5.2), Inches(0.6),
                 font_size=12, color=COLOR_DARK_GRAY)

# Win-Win section
add_rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(2.05), fill_color=COLOR_WHITE)
add_rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(0.4), fill_color=COLOR_NAVY)
add_text_box(slide, "両社の協業により実現できること", Inches(0.4), Inches(5.1), Inches(12.5), Inches(0.4),
             font_size=14, bold=True, color=COLOR_WHITE)

win_items = [
    ("CCT様", "FM/建築DX領域という新たな収益源の獲得", COLOR_ACCENT),
    ("梓総研", "プロダクトの市場浸透と技術開発への集中", COLOR_GREEN),
    ("エンドユーザー", "ワンストップ（Notion＋AIR-Plate＋導入支援）での価値提供", COLOR_ORANGE),
]
for i, (who, what, color) in enumerate(win_items):
    x = Inches(0.5) + i * Inches(4.2)
    add_rect(slide, x, Inches(5.6), Inches(3.9), Inches(1.4), fill_color=COLOR_LIGHT_GRAY)
    add_rect(slide, x, Inches(5.6), Inches(3.9), Inches(0.3), fill_color=color)
    add_text_box(slide, who, x + Inches(0.1), Inches(5.6), Inches(3.7), Inches(0.3),
                 font_size=11, bold=True, color=COLOR_WHITE)
    add_text_box(slide, what, x + Inches(0.1), Inches(5.95), Inches(3.7), Inches(0.95),
                 font_size=12, color=COLOR_DARK_GRAY)

# Right side highlight box
add_rect(slide, Inches(8.3), Inches(1.4), Inches(4.7), Inches(3.5), fill_color=COLOR_LIGHT_BLUE)
add_rect(slide, Inches(8.3), Inches(1.4), Inches(0.07), Inches(3.5), fill_color=COLOR_ORANGE)
add_text_box(slide, "なぜ今、CCT様なのか",
             Inches(8.5), Inches(1.5), Inches(4.4), Inches(0.5),
             font_size=14, bold=True, color=COLOR_NAVY)
add_text_box(slide,
             "Notionを基盤とするAIR-Plateにとって、\nNotion代理店としての実績と顧客基盤を\n持つCCT様は最高のパートナーです。\n\n"
             "建築DX × IT × ビジネス開発の\n三位一体の強みで、市場の新標準を\n共に創りましょう。",
             Inches(8.5), Inches(2.1), Inches(4.3), Inches(2.6),
             font_size=13, color=COLOR_NAVY)


# ─────────────────────────────────────────
# SLIDE 5: Collaboration Scheme
# ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
add_rect(slide, 0, 0, Inches(13.33), Inches(7.5), fill_color=RGBColor(0xF7, 0xF9, 0xFC))
add_slide_header(slide, "協業スキーム ― 役割分担の全体像",
                 "CCT様が「事業主体」として自社名で契約・請求できるスキーム")

# Table header
table_top = Inches(1.45)
col_widths = [Inches(2.2), Inches(4.9), Inches(5.7)]
col_x = [Inches(0.3), Inches(2.5), Inches(7.4)]

headers = ["項目", "梓総研（甲）", "CCT様（乙）"]
header_colors = [COLOR_NAVY, COLOR_NAVY, COLOR_ACCENT]
for i, (hdr, col, cx) in enumerate(zip(headers, header_colors, col_x)):
    add_rect(slide, cx, table_top, col_widths[i], Inches(0.45), fill_color=col)
    add_text_box(slide, hdr, cx + Inches(0.1), table_top, col_widths[i], Inches(0.45),
                 font_size=14, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("知財・開発", "AIR-Plate本体の知財保持\nテンプレート継続アップデート", "―"),
    ("技術支援", "設計仕様の開示、技術レクチャー\n二次サポート", "一次サポート・導入支援\n保守対応"),
    ("営業・販売", "案件レジストレーションの承認", "事業主体として販売・マーケティング\n契約締結・請求"),
    ("エンドユーザー対応", "建築専門知識を要する\n高度な疑義対応", "デモ・要件定義・設定支援\n運用サポート全般"),
]
row_colors = [COLOR_LIGHT_GRAY, COLOR_WHITE, COLOR_LIGHT_GRAY, COLOR_WHITE]
row_label_colors = [COLOR_NAVY, COLOR_ACCENT, COLOR_NAVY, COLOR_ACCENT]

for r, (label, azusa, cct) in enumerate(rows):
    y = table_top + Inches(0.45) + r * Inches(1.1)
    bg = row_colors[r]
    # Row background
    for i, (cx, cw) in enumerate(zip(col_x, col_widths)):
        add_rect(slide, cx, y, cw, Inches(1.1), fill_color=bg)
    # Label
    add_rect(slide, col_x[0], y, col_widths[0], Inches(1.1), fill_color=row_label_colors[r])
    add_text_box(slide, label, col_x[0] + Inches(0.1), y + Inches(0.25),
                 col_widths[0] - Inches(0.2), Inches(0.6),
                 font_size=13, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    # Azusa content
    add_text_box(slide, azusa, col_x[1] + Inches(0.15), y + Inches(0.15),
                 col_widths[1] - Inches(0.3), Inches(0.9), font_size=12, color=COLOR_DARK_GRAY)
    # CCT content
    add_text_box(slide, cct, col_x[2] + Inches(0.15), y + Inches(0.15),
                 col_widths[2] - Inches(0.3), Inches(0.9), font_size=12, color=COLOR_DARK_GRAY)

# Highlight
add_rect(slide, Inches(0.3), Inches(6.1), Inches(12.7), Inches(0.4),
         fill_color=COLOR_ORANGE)
add_text_box(slide,
             "ポイント：CCT様が「事業主体」として自社名で契約・請求できるスキームです",
             Inches(0.5), Inches(6.1), Inches(12.5), Inches(0.4),
             font_size=13, bold=True, color=COLOR_WHITE)


# ─────────────────────────────────────────
# SLIDE 6: CCT Business Image
# ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
add_rect(slide, 0, 0, Inches(13.33), Inches(7.5), fill_color=RGBColor(0xF7, 0xF9, 0xFC))
add_slide_header(slide, "CCT様の業務イメージ", "販売代理店としての主な業務")

tasks = [
    ("1", "営業・提案活動",
     "エンドユーザーへのAIR-Plate＋Notionのセット提案\nデモンストレーション・PoC支援", COLOR_ACCENT),
    ("2", "契約・請求の一本化",
     "Notionライセンス＋AIR-Plateライセンスをセットで請求\nエンドユーザーとの利用契約はCCT様名義で締結", COLOR_NAVY),
    ("3", "導入・運用サポート",
     "要件定義・初期設定・カスタマイズ\n運用開始後の問い合わせ一次対応", COLOR_GREEN),
    ("4", "独自付加価値の提供",
     "CCT様のコンサルティングサービスとの組み合わせ提案も可能", COLOR_ORANGE),
]

for i, (num, title, desc, color) in enumerate(tasks):
    col = i % 2
    row = i // 2
    x = Inches(0.3) + col * Inches(6.4)
    y = Inches(1.5) + row * Inches(2.8)
    w = Inches(6.2)
    h = Inches(2.55)

    add_rect(slide, x, y, w, h, fill_color=COLOR_WHITE)
    add_rect(slide, x, y, w, Inches(0.5), fill_color=color)

    # Number badge
    add_rect(slide, x + Inches(0.15), y + Inches(0.08), Inches(0.34), Inches(0.34),
             fill_color=COLOR_WHITE)
    add_text_box(slide, num, x + Inches(0.15), y + Inches(0.08), Inches(0.34), Inches(0.34),
                 font_size=14, bold=True, color=color, align=PP_ALIGN.CENTER)

    add_text_box(slide, title, x + Inches(0.6), y + Inches(0.08), w - Inches(0.7), Inches(0.38),
                 font_size=14, bold=True, color=COLOR_WHITE)
    add_text_box(slide, desc, x + Inches(0.2), y + Inches(0.65), w - Inches(0.4), Inches(1.75),
                 font_size=13, color=COLOR_DARK_GRAY)


# ─────────────────────────────────────────
# SLIDE 7: Enablement Support
# ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
add_rect(slide, 0, 0, Inches(13.33), Inches(7.5), fill_color=RGBColor(0xF7, 0xF9, 0xFC))
add_slide_header(slide, "梓総研が提供するイネーブルメント支援",
                 "CCT様が自律的に販売・導入できるよう、梓総研が全面サポート")

supports = [
    ("技術開示", "テンプレートの設計仕様・数式・リレーション\n構造等の詳細資料を提供", COLOR_ACCENT),
    ("定期レクチャー", "アップデート内容の共有、\n新機能の使い方研修", COLOR_NAVY),
    ("担当者トレーニング", "CCT様の営業・技術担当者向け\nハンズオン研修", COLOR_GREEN),
    ("二次サポート", "建築専門知識を要する高度な\n問い合わせへのバックアップ", COLOR_ORANGE),
]

for i, (title, desc, color) in enumerate(supports):
    x = Inches(0.3) + i * Inches(3.2)
    y = Inches(1.5)
    w = Inches(3.0)
    h = Inches(4.5)

    add_rect(slide, x, y, w, h, fill_color=COLOR_WHITE)
    # Top color strip
    add_rect(slide, x, y, w, Inches(0.5), fill_color=color)
    # Icon placeholder circle
    cx_c = x + w/2 - Inches(0.4)
    add_rect(slide, cx_c, y + Inches(0.65), Inches(0.8), Inches(0.8), fill_color=COLOR_LIGHT_BLUE)

    add_text_box(slide, title, x, y + Inches(0.1), w, Inches(0.38),
                 font_size=14, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, x + Inches(0.15), y + Inches(1.6), w - Inches(0.3), Inches(2.6),
                 font_size=13, color=COLOR_DARK_GRAY, align=PP_ALIGN.CENTER)

# Bottom note
add_rect(slide, Inches(0.3), Inches(6.2), Inches(12.7), Inches(0.95), fill_color=COLOR_LIGHT_BLUE)
add_rect(slide, Inches(0.3), Inches(6.2), Inches(0.07), Inches(0.95), fill_color=COLOR_ACCENT)
add_text_box(slide,
             "→ 段階的にCCT様の独力対応範囲を広げていくイメージです。\n"
             "最初は梓総研が手厚くサポートし、実績を積みながらCCT様が自立的に展開できる体制を構築します。",
             Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.8),
             font_size=12, color=COLOR_NAVY, italic=True)


# ─────────────────────────────────────────
# SLIDE 8: Revenue Model
# ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
add_rect(slide, 0, 0, Inches(13.33), Inches(7.5), fill_color=RGBColor(0xF7, 0xF9, 0xFC))
add_slide_header(slide, "収益モデル", "ロイヤリティ還元方式 ― CCT様が最大限の利益を享受できる構造")

# Revenue flow diagram
add_text_box(slide, "収益フロー", Inches(0.3), Inches(1.4), Inches(12.7), Inches(0.4),
             font_size=14, bold=True, color=COLOR_NAVY)

# Box: End user
add_rect(slide, Inches(0.3), Inches(2.0), Inches(2.8), Inches(1.2), fill_color=COLOR_LIGHT_GRAY)
add_rect(slide, Inches(0.3), Inches(2.0), Inches(2.8), Inches(0.35), fill_color=COLOR_MID_GRAY)
add_text_box(slide, "エンドユーザー", Inches(0.35), Inches(2.0), Inches(2.7), Inches(0.35),
             font_size=12, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, "ライセンス料＋\n導入費用を支払い", Inches(0.4), Inches(2.4), Inches(2.6), Inches(0.7),
             font_size=11, color=COLOR_DARK_GRAY, align=PP_ALIGN.CENTER)

# Arrow
add_rect(slide, Inches(3.2), Inches(2.45), Inches(0.8), Inches(0.1), fill_color=COLOR_ACCENT)
add_text_box(slide, "→", Inches(3.2), Inches(2.3), Inches(0.8), Inches(0.4),
             font_size=18, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)

# Box: CCT
add_rect(slide, Inches(4.1), Inches(2.0), Inches(3.5), Inches(1.2), fill_color=COLOR_LIGHT_BLUE)
add_rect(slide, Inches(4.1), Inches(2.0), Inches(3.5), Inches(0.35), fill_color=COLOR_ACCENT)
add_text_box(slide, "CCT様（乙）", Inches(4.15), Inches(2.0), Inches(3.4), Inches(0.35),
             font_size=12, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, "売上を受領\n（70〜80%を取得）", Inches(4.2), Inches(2.4), Inches(3.2), Inches(0.7),
             font_size=11, color=COLOR_ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Arrow
add_text_box(slide, "→", Inches(7.7), Inches(2.3), Inches(0.8), Inches(0.4),
             font_size=18, bold=True, color=COLOR_ORANGE, align=PP_ALIGN.CENTER)
add_text_box(slide, "ロイヤリティ\n(20〜30%)", Inches(7.6), Inches(2.7), Inches(1.0), Inches(0.55),
             font_size=9, color=COLOR_ORANGE, align=PP_ALIGN.CENTER)

# Box: Azusa
add_rect(slide, Inches(8.6), Inches(2.0), Inches(3.5), Inches(1.2), fill_color=COLOR_LIGHT_GRAY)
add_rect(slide, Inches(8.6), Inches(2.0), Inches(3.5), Inches(0.35), fill_color=COLOR_NAVY)
add_text_box(slide, "梓総研（甲）", Inches(8.65), Inches(2.0), Inches(3.4), Inches(0.35),
             font_size=12, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, "ロイヤリティを受領\n（AIR-Plate相当分の20〜30%）", Inches(8.7), Inches(2.4), Inches(3.2), Inches(0.7),
             font_size=11, color=COLOR_DARK_GRAY, align=PP_ALIGN.CENTER)

# Royalty rate details
add_rect(slide, Inches(0.3), Inches(3.5), Inches(6.1), Inches(3.1), fill_color=COLOR_WHITE)
add_rect(slide, Inches(0.3), Inches(3.5), Inches(6.1), Inches(0.4), fill_color=COLOR_ACCENT)
add_text_box(slide, "ロイヤリティ料率（ご協議事項）", Inches(0.4), Inches(3.5), Inches(5.9), Inches(0.4),
             font_size=13, bold=True, color=COLOR_WHITE)

royalty_data = [
    ("想定レンジ", "20〜30%程度"),
    ("CCT様の取り分", "70〜80%"),
    ("算定基礎", "AIR-Plate相当分のライセンス収益"),
]
for i, (k, v) in enumerate(royalty_data):
    y = Inches(4.0) + i * Inches(0.65)
    add_rect(slide, Inches(0.45), y, Inches(2.2), Inches(0.55), fill_color=COLOR_LIGHT_BLUE)
    add_text_box(slide, k, Inches(0.5), y + Inches(0.07), Inches(2.1), Inches(0.42),
                 font_size=12, bold=True, color=COLOR_NAVY, align=PP_ALIGN.CENTER)
    add_text_box(slide, v, Inches(2.75), y + Inches(0.07), Inches(3.5), Inches(0.42),
                 font_size=13, bold=True, color=COLOR_ACCENT)

# Individual support fees
add_rect(slide, Inches(6.6), Inches(3.5), Inches(6.4), Inches(3.1), fill_color=COLOR_WHITE)
add_rect(slide, Inches(6.6), Inches(3.5), Inches(6.4), Inches(0.4), fill_color=COLOR_ORANGE)
add_text_box(slide, "個別技術支援費用", Inches(6.7), Inches(3.5), Inches(6.2), Inches(0.4),
             font_size=13, bold=True, color=COLOR_WHITE)
add_text_box(slide,
             "イネーブルメント支援の範囲を超える\n案件固有のカスタマイズ・コンサルティングは、\n別途出来高ベースで協議。\n\n"
             "※ 具体的な料率・算定基礎は\n　 本日のご意見を踏まえて詰めてまいります。",
             Inches(6.75), Inches(4.0), Inches(6.1), Inches(2.4),
             font_size=12, color=COLOR_DARK_GRAY, italic=False)


# ─────────────────────────────────────────
# SLIDE 9: Registration System
# ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
add_rect(slide, 0, 0, Inches(13.33), Inches(7.5), fill_color=RGBColor(0xF7, 0xF9, 0xFC))
add_slide_header(slide, "案件の進め方 ― レジストレーション制度",
                 "双方にとって公平・透明な案件管理の仕組み")

# Registration flow
add_rect(slide, Inches(0.3), Inches(1.45), Inches(12.7), Inches(2.9), fill_color=COLOR_WHITE)
add_rect(slide, Inches(0.3), Inches(1.45), Inches(12.7), Inches(0.4), fill_color=COLOR_NAVY)
add_text_box(slide, "案件登録（レジストレーション）の仕組み", Inches(0.4), Inches(1.45), Inches(12.5), Inches(0.4),
             font_size=14, bold=True, color=COLOR_WHITE)

steps = [
    ("STEP 1", "CCT様が\n営業活動開始前に、\n対象案件を梓総研へ登録"),
    ("STEP 2", "梓総研が承認\n↓\nCCT様の優先権を保護"),
    ("STEP 3", "承認済み案件には\n梓総研は\n直接営業を行わない"),
]
step_colors = [COLOR_ACCENT, COLOR_NAVY, COLOR_GREEN]
for i, (step, desc) in enumerate(steps):
    x = Inches(0.5) + i * Inches(4.1)
    y = Inches(2.0)
    add_rect(slide, x, y, Inches(3.7), Inches(0.4), fill_color=step_colors[i])
    add_text_box(slide, step, x, y, Inches(3.7), Inches(0.4),
                 font_size=13, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x, y + Inches(0.4), Inches(3.7), Inches(1.4), fill_color=COLOR_LIGHT_GRAY)
    add_text_box(slide, desc, x + Inches(0.1), y + Inches(0.5), Inches(3.5), Inches(1.25),
                 font_size=13, color=COLOR_DARK_GRAY, align=PP_ALIGN.CENTER)
    if i < 2:
        add_text_box(slide, "→", Inches(4.25) + i * Inches(4.1), y + Inches(0.7), Inches(0.5), Inches(0.4),
                     font_size=20, bold=True, color=step_colors[i+1], align=PP_ALIGN.CENTER)

# Case separation
add_rect(slide, Inches(0.3), Inches(4.55), Inches(12.7), Inches(2.6), fill_color=COLOR_WHITE)
add_rect(slide, Inches(0.3), Inches(4.55), Inches(12.7), Inches(0.4), fill_color=COLOR_ACCENT)
add_text_box(slide, "梓設計案件との区分け", Inches(0.4), Inches(4.55), Inches(12.5), Inches(0.4),
             font_size=14, bold=True, color=COLOR_WHITE)

cases = [
    ("梓設計グループが主体の既存案件",
     "梓設計から業務委託として参画いただく形",
     COLOR_NAVY, "支援スタッフとして入るケース"),
    ("CCT様主導の新規案件",
     "代理店として自社名義で推進",
     COLOR_ACCENT, "代理店として売るケース"),
]
for i, (situation, action, color, label) in enumerate(cases):
    x = Inches(0.5) + i * Inches(6.3)
    y = Inches(5.1)
    add_rect(slide, x, y, Inches(5.9), Inches(1.8), fill_color=COLOR_LIGHT_GRAY)
    add_rect(slide, x, y, Inches(5.9), Inches(0.35), fill_color=color)
    add_text_box(slide, label, x, y, Inches(5.9), Inches(0.35),
                 font_size=11, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, f"状況：{situation}", x + Inches(0.15), y + Inches(0.45), Inches(5.6), Inches(0.5),
                 font_size=12, bold=True, color=COLOR_NAVY)
    add_text_box(slide, f"対応：{action}", x + Inches(0.15), y + Inches(0.95), Inches(5.6), Inches(0.5),
                 font_size=12, color=COLOR_DARK_GRAY)


# ─────────────────────────────────────────
# SLIDE 10: Project Pipeline
# ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
add_rect(slide, 0, 0, Inches(13.33), Inches(7.5), fill_color=RGBColor(0xF7, 0xF9, 0xFC))
add_slide_header(slide, "案件の整理（短期・長期）",
                 "段階的な協業展開で確実にパートナーシップを構築")

# Short-term
add_rect(slide, Inches(0.3), Inches(1.45), Inches(6.1), Inches(4.0), fill_color=COLOR_WHITE)
add_rect(slide, Inches(0.3), Inches(1.45), Inches(6.1), Inches(0.5), fill_color=COLOR_ACCENT)
add_text_box(slide, "短期：協業試行案件", Inches(0.4), Inches(1.45), Inches(5.9), Inches(0.5),
             font_size=15, bold=True, color=COLOR_WHITE)
add_text_box(slide, "まずはここから連携を開始", Inches(0.4), Inches(2.05), Inches(5.9), Inches(0.35),
             font_size=11, color=COLOR_ACCENT, italic=True)

short_term = ["羽田空港", "青山葬儀所", "バンテリンドーム", "三菱電機"]
for i, project in enumerate(short_term):
    y = Inches(2.5) + i * Inches(0.65)
    add_rect(slide, Inches(0.5), y, Inches(0.3), Inches(0.3), fill_color=COLOR_ACCENT)
    add_text_box(slide, project, Inches(0.95), y - Inches(0.02), Inches(5.2), Inches(0.38),
                 font_size=14, color=COLOR_DARK_GRAY)

# Long-term
add_rect(slide, Inches(6.6), Inches(1.45), Inches(6.4), Inches(4.0), fill_color=COLOR_WHITE)
add_rect(slide, Inches(6.6), Inches(1.45), Inches(6.4), Inches(0.5), fill_color=COLOR_NAVY)
add_text_box(slide, "長期：展開候補案件", Inches(6.7), Inches(1.45), Inches(6.2), Inches(0.5),
             font_size=15, bold=True, color=COLOR_WHITE)
add_text_box(slide, "実績を積んだ上で順次拡大", Inches(6.7), Inches(2.05), Inches(6.2), Inches(0.35),
             font_size=11, color=COLOR_NAVY, italic=True)

long_term = ["成田空港", "ダイワロイネット", "羽田空港（拡張）", "竹中工務店全般"]
for i, project in enumerate(long_term):
    y = Inches(2.5) + i * Inches(0.65)
    add_rect(slide, Inches(6.8), y, Inches(0.3), Inches(0.3), fill_color=COLOR_NAVY)
    add_text_box(slide, project, Inches(7.25), y - Inches(0.02), Inches(5.6), Inches(0.38),
                 font_size=14, color=COLOR_DARK_GRAY)

# Arrow between
add_text_box(slide, "→\n実績\n拡大", Inches(6.05), Inches(2.7), Inches(0.55), Inches(1.5),
             font_size=11, bold=True, color=COLOR_ORANGE, align=PP_ALIGN.CENTER)

# Bottom approach
add_rect(slide, Inches(0.3), Inches(5.6), Inches(12.7), Inches(1.55), fill_color=COLOR_LIGHT_BLUE)
add_rect(slide, Inches(0.3), Inches(5.6), Inches(0.07), Inches(1.55), fill_color=COLOR_ORANGE)
add_text_box(slide, "進め方",
             Inches(0.5), Inches(5.65), Inches(12.3), Inches(0.4),
             font_size=13, bold=True, color=COLOR_NAVY)
add_text_box(slide,
             "まずは短期案件で協業を試行し、業務フロー・サポート体制を検証。\n"
             "実績を積みながら、段階的に長期案件・大規模案件へ拡大。",
             Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.9),
             font_size=13, color=COLOR_DARK_GRAY)


# ─────────────────────────────────────────
# SLIDE 11: IP & NDA
# ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
add_rect(slide, 0, 0, Inches(13.33), Inches(7.5), fill_color=RGBColor(0xF7, 0xF9, 0xFC))
add_slide_header(slide, "知的財産権・秘密保持",
                 "双方のノウハウを守り、プロダクト全体の品質向上に貢献する仕組み")

# IP section
add_rect(slide, Inches(0.3), Inches(1.45), Inches(6.1), Inches(3.5), fill_color=COLOR_WHITE)
add_rect(slide, Inches(0.3), Inches(1.45), Inches(6.1), Inches(0.45), fill_color=COLOR_NAVY)
add_text_box(slide, "知的財産権", Inches(0.4), Inches(1.45), Inches(5.9), Inches(0.45),
             font_size=15, bold=True, color=COLOR_WHITE)

ip_items = [
    "AIR-Plate本体・標準テンプレートの\n知財は梓総研に帰属",
    "CCT様がエンドユーザー向けに作成した設定・追加機能のうち、\n汎用性のあるものは協議の上、次期アップデートへ反映可能",
]
for i, item in enumerate(ip_items):
    y = Inches(2.05) + i * Inches(1.3)
    add_rect(slide, Inches(0.5), y, Inches(0.35), Inches(0.35), fill_color=COLOR_NAVY)
    add_text_box(slide, str(i+1), Inches(0.5), y, Inches(0.35), Inches(0.35),
                 font_size=12, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, item, Inches(0.95), y - Inches(0.05), Inches(5.2), Inches(1.1),
                 font_size=13, color=COLOR_DARK_GRAY)

# NDA section
add_rect(slide, Inches(6.6), Inches(1.45), Inches(6.4), Inches(3.5), fill_color=COLOR_WHITE)
add_rect(slide, Inches(6.6), Inches(1.45), Inches(6.4), Inches(0.45), fill_color=COLOR_ACCENT)
add_text_box(slide, "秘密保持", Inches(6.7), Inches(1.45), Inches(6.2), Inches(0.45),
             font_size=15, bold=True, color=COLOR_WHITE)

nda_items = [
    "双方が開示する技術・営業情報について、\n相互に守秘義務を負う",
    "第三者への開示は\n事前の書面承諾が必要",
]
for i, item in enumerate(nda_items):
    y = Inches(2.05) + i * Inches(1.3)
    add_rect(slide, Inches(6.8), y, Inches(0.35), Inches(0.35), fill_color=COLOR_ACCENT)
    add_text_box(slide, str(i+1), Inches(6.8), y, Inches(0.35), Inches(0.35),
                 font_size=12, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, item, Inches(7.25), y - Inches(0.05), Inches(5.6), Inches(1.1),
                 font_size=13, color=COLOR_DARK_GRAY)

# Summary box
add_rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(2.05), fill_color=COLOR_LIGHT_BLUE)
add_rect(slide, Inches(0.3), Inches(5.1), Inches(0.07), Inches(2.05), fill_color=COLOR_ORANGE)
add_text_box(slide,
             "→ CCT様の独自ノウハウを尊重しつつ、プロダクト全体の品質向上にも資する仕組みです。\n\n"
             "CCT様が現場で培った知見がAIR-Plateをより良くし、より多くのユーザーに価値を届けます。",
             Inches(0.5), Inches(5.15), Inches(12.3), Inches(1.85),
             font_size=13, color=COLOR_NAVY)


# ─────────────────────────────────────────
# SLIDE 12: Contract Terms
# ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
add_rect(slide, 0, 0, Inches(13.33), Inches(7.5), fill_color=RGBColor(0xF7, 0xF9, 0xFC))
add_slide_header(slide, "契約の基本条件", "販売代理店契約の主要条件")

contract_items = [
    ("契約形態", "販売代理店契約", COLOR_NAVY),
    ("有効期間", "締結日から1年間（自動更新）", COLOR_ACCENT),
    ("更新拒絶", "期間満了3ヶ月前までに書面通知", COLOR_GREEN),
    ("対価", "ロイヤリティ＋個別技術支援費（出来高ベース）", COLOR_ORANGE),
]

for i, (label, value, color) in enumerate(contract_items):
    y = Inches(1.6) + i * Inches(1.2)
    # Full row background
    add_rect(slide, Inches(0.3), y, Inches(12.7), Inches(1.0), fill_color=COLOR_WHITE)
    # Label
    add_rect(slide, Inches(0.3), y, Inches(2.8), Inches(1.0), fill_color=color)
    add_text_box(slide, label, Inches(0.35), y + Inches(0.2), Inches(2.7), Inches(0.6),
                 font_size=16, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    # Value
    add_text_box(slide, value, Inches(3.3), y + Inches(0.22), Inches(9.5), Inches(0.55),
                 font_size=16, color=COLOR_DARK_GRAY)
    # Accent dot
    add_rect(slide, Inches(3.1), y + Inches(0.38), Inches(0.12), Inches(0.12), fill_color=color)

# Note
add_rect(slide, Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.6), fill_color=COLOR_LIGHT_BLUE)
add_text_box(slide,
             "※ 上記は基本条件の骨格です。詳細条件は本日の協議内容を踏まえ、契約書ドラフトに反映いたします。",
             Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
             font_size=12, color=COLOR_NAVY, italic=True)


# ─────────────────────────────────────────
# SLIDE 13: Confirmation Points
# ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
add_rect(slide, 0, 0, Inches(13.33), Inches(7.5), fill_color=RGBColor(0xF7, 0xF9, 0xFC))
add_slide_header(slide, "本日ご確認いただきたいポイント",
                 "以下の点について、CCT様のお考え・ご要望をお聞かせください")

confirm_items = [
    ("1", "役割分担について", "業務範囲の認識にずれがないか"),
    ("2", "ロイヤリティ料率について", "想定レンジ（20〜30%）に対するご意見"),
    ("3", "ロイヤリティ算定基礎について", "AIR-Plate相当分 or セット価格全体"),
    ("4", "案件レジストレーション制度について", "運用上のご懸念"),
    ("5", "梓設計案件との区分けについて", "代理店 vs 業務委託の使い分け"),
    ("6", "短期案件での試行について", "開始時期・対象案件のご希望"),
]

col1 = confirm_items[:3]
col2 = confirm_items[3:]

for col_idx, items in enumerate([col1, col2]):
    x = Inches(0.3) + col_idx * Inches(6.5)
    for i, (num, title, question) in enumerate(items):
        y = Inches(1.55) + i * Inches(1.7)
        add_rect(slide, x, y, Inches(6.2), Inches(1.55), fill_color=COLOR_WHITE)
        add_rect(slide, x, y, Inches(0.5), Inches(1.55),
                 fill_color=COLOR_ACCENT if col_idx == 0 else COLOR_NAVY)
        add_text_box(slide, num, x, y + Inches(0.55), Inches(0.5), Inches(0.45),
                     font_size=16, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, title, x + Inches(0.6), y + Inches(0.15), Inches(5.5), Inches(0.5),
                     font_size=13, bold=True, color=COLOR_NAVY)
        add_text_box(slide, f"▶ {question}", x + Inches(0.6), y + Inches(0.7), Inches(5.5), Inches(0.6),
                     font_size=12, color=COLOR_MID_GRAY)


# ─────────────────────────────────────────
# SLIDE 14: Schedule
# ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
add_rect(slide, 0, 0, Inches(13.33), Inches(7.5), fill_color=RGBColor(0xF7, 0xF9, 0xFC))
add_slide_header(slide, "今後のスケジュール（案）", "スピーディーな合意形成・契約締結を目指します")

schedule = [
    ("本日", "協業フレームワークの方向性合意", COLOR_ACCENT, "Step 1"),
    ("〜2週間後", "ロイヤリティ料率・算定基礎の合意\n契約書ドラフト修正", COLOR_NAVY, "Step 2"),
    ("〜1ヶ月後", "契約書最終確認・締結", COLOR_GREEN, "Step 3"),
    ("契約締結後", "イネーブルメント支援開始\n短期案件での協業試行", COLOR_ORANGE, "Step 4"),
]

for i, (timing, content, color, step) in enumerate(schedule):
    y = Inches(1.6) + i * Inches(1.35)
    # Step badge
    add_rect(slide, Inches(0.3), y, Inches(1.2), Inches(1.1), fill_color=color)
    add_text_box(slide, step, Inches(0.3), y + Inches(0.05), Inches(1.2), Inches(0.35),
                 font_size=10, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, timing, Inches(0.3), y + Inches(0.45), Inches(1.2), Inches(0.55),
                 font_size=13, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    # Content
    add_rect(slide, Inches(1.5), y, Inches(11.5), Inches(1.1),
             fill_color=COLOR_WHITE if i % 2 == 0 else COLOR_LIGHT_GRAY)
    add_rect(slide, Inches(1.5), y, Inches(0.05), Inches(1.1), fill_color=color)
    add_text_box(slide, content, Inches(1.7), y + Inches(0.15), Inches(11.1), Inches(0.85),
                 font_size=15, color=COLOR_DARK_GRAY)


# ─────────────────────────────────────────
# SLIDE 15: Closing
# ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])

add_rect(slide, 0, 0, Inches(13.33), Inches(7.5), fill_color=COLOR_NAVY)
add_rect(slide, 0, 0, Inches(13.33), Inches(2.5), fill_color=RGBColor(0x0D, 0x1E, 0x4A))
add_rect(slide, Inches(10.0), 0, Inches(3.33), Inches(7.5), fill_color=COLOR_ACCENT)
add_rect(slide, Inches(9.3), 0, Inches(0.8), Inches(7.5), fill_color=RGBColor(0x00, 0x5E, 0x9A))
add_rect(slide, 0, Inches(2.5), Inches(13.33), Inches(0.04), fill_color=COLOR_ORANGE)

add_text_box(slide, "AIR-Plate  ×  CCT", Inches(0.5), Inches(0.6), Inches(8.5), Inches(1.0),
             font_size=40, bold=True, color=COLOR_WHITE, align=PP_ALIGN.LEFT)
add_text_box(slide, "FM / 建築DXの新たなスタンダードへ",
             Inches(0.5), Inches(1.55), Inches(8.5), Inches(0.65),
             font_size=20, bold=False, color=RGBColor(0xB0, 0xC8, 0xE8), align=PP_ALIGN.LEFT)

add_text_box(slide, "梓総研の建築専門知識  ×  CCT様のIT力・営業力",
             Inches(0.5), Inches(3.0), Inches(8.5), Inches(0.5),
             font_size=16, bold=True, color=COLOR_WHITE)
add_text_box(slide, "双方の強みを活かした、Win-Winの事業パートナーシップ",
             Inches(0.5), Inches(3.55), Inches(8.5), Inches(0.45),
             font_size=14, color=RGBColor(0xB0, 0xC8, 0xE8))

add_rect(slide, Inches(0.5), Inches(4.1), Inches(5.0), Inches(0.03), fill_color=COLOR_ORANGE)

add_text_box(slide,
             "本日の議論を踏まえ、両社にとって最良のスキームを\n一緒に作り上げてまいりたいと考えております。\n\n忌憚のないご意見をお聞かせください。",
             Inches(0.5), Inches(4.3), Inches(8.5), Inches(2.0),
             font_size=15, color=COLOR_WHITE)

add_text_box(slide, "ありがとうございました", Inches(0.5), Inches(6.6), Inches(8.5), Inches(0.5),
             font_size=18, bold=True, color=RGBColor(0xB0, 0xC8, 0xE8))

add_text_box(slide, "株式会社梓総合研究所\n2026年5月", Inches(10.2), Inches(2.8), Inches(2.8), Inches(1.0),
             font_size=13, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

# Save
output_path = "/home/user/Claude_Test/AIRPlate_CCT_Partnership_Proposal.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
