from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── カラーパレット（淡いトーン）────────────────────
C_ACCENT  = RGBColor(0x4A, 0x7C, 0xB4)   # ミディアムブルー（線・アクセント）
C_LIGHT   = RGBColor(0xE4, 0xEF, 0xF7)   # 薄青（テーブルヘッダ・帯）
C_PALE    = RGBColor(0xF4, 0xF7, 0xFB)   # 極薄青（交互行・ボックス背景）
C_GRAY    = RGBColor(0xF5, 0xF5, 0xF5)   # 薄グレー（交互行）
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK    = RGBColor(0x22, 0x22, 0x22)   # ほぼ黒
C_MID     = RGBColor(0x4A, 0x7C, 0xB4)   # 見出し文字
C_SUB     = RGBColor(0x55, 0x55, 0x55)   # サブテキスト
C_WARN    = RGBColor(0x8B, 0x60, 0x00)   # 注釈（茶）
C_AMBER   = RGBColor(0xFF, 0xF3, 0xDC)   # 議題ボックス背景（淡黄）
C_AMBER_B = RGBColor(0xE0, 0xA0, 0x30)   # 議題ボックスボーダー

# ── ページサイズ（A4縦）──────────────────────────
W = Inches(8.27)   # A4幅
H = Inches(11.69)  # A4高さ

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ══ ユーティリティ ════════════════════════════════

def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_pt=1.0):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=14, bold=False, color=C_DARK,
             align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def section_title(slide, text, top):
    """細い青線＋青文字の見出し"""
    add_rect(slide, Inches(0.5), top, Inches(7.27), Pt(2),
             fill=C_ACCENT)
    add_text(slide, text,
             Inches(0.5), top + Pt(4), Inches(7.27), Inches(0.45),
             size=17, bold=True, color=C_MID)
    return top + Inches(0.52)

def sub_heading(slide, text, top, left=Inches(0.5)):
    """小見出し（太字・アクセントカラー）"""
    add_text(slide, text, left, top, Inches(7.27), Inches(0.38),
             size=14, bold=True, color=C_MID)
    return top + Inches(0.42)

def page_footer(slide, num, total=7):
    add_text(slide,
             "コンフォール井の頭公園南 修繕委員会 第1回定例　2026.6.7",
             Inches(0.5), Inches(11.2), Inches(6.5), Inches(0.3),
             size=10, color=C_SUB)
    add_text(slide, f"{num} / {total}",
             Inches(7.5), Inches(11.2), Inches(0.6), Inches(0.3),
             size=10, color=C_SUB, align=PP_ALIGN.RIGHT)
    # 下線
    add_rect(slide, Inches(0.5), Inches(11.15), Inches(7.27), Pt(1),
             fill=C_LIGHT)

def add_table(slide, headers, rows, l, t, w, h,
              hdr_size=13, data_size=13, col_widths=None):
    cols = len(headers)
    tbl = slide.shapes.add_table(len(rows)+1, cols, l, t, w, h).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    else:
        cw = w // cols
        for i in range(cols):
            tbl.columns[i].width = cw

    for ci, hd in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_LIGHT
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.text = hd
        run.font.size = Pt(hdr_size)
        run.font.bold = True
        run.font.color.rgb = C_MID

    for ri, row in enumerate(rows):
        bg = C_PALE if ri % 2 == 0 else C_WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri+1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.text = str(val)
            run.font.size = Pt(data_size)
            run.font.color.rgb = C_DARK
    return tbl

def note_box(slide, text, l, t, w, h, bg=C_AMBER, border=C_AMBER_B):
    add_rect(slide, l, t, w, h, fill=bg, line_color=border, line_pt=0.75)
    add_text(slide, text, l+Inches(0.12), t+Inches(0.06),
             w-Inches(0.18), h-Inches(0.1),
             size=13, color=RGBColor(0x6B,0x40,0x00))

def bullet(slide, items, l, t, w, size=14, gap=None):
    """箇条書きを縦に並べて返す最終y座標"""
    if gap is None:
        gap = Inches(0.48) if size >= 14 else Inches(0.42)
    y = t
    for item in items:
        add_text(slide, item, l, y, w, gap + Inches(0.05),
                 size=size, color=C_DARK)
        y += gap
    return y

# ══════════════════════════════════════════════════════
# PAGE 1 ｜ 表紙
# ══════════════════════════════════════════════════════
p1 = prs.slides.add_slide(BLANK)

# 上部アクセントバー（細め）
add_rect(p1, 0, 0, W, Inches(0.18), fill=C_ACCENT)

# タイトルエリア（白背景・薄青の枠）
add_rect(p1, Inches(0.6), Inches(2.8), Inches(7.07), Inches(3.5),
         fill=C_PALE, line_color=C_ACCENT, line_pt=1.5)

add_text(p1, "コンフォール井の頭公園南", Inches(0.6), Inches(3.0), Inches(7.07), Inches(0.6),
         size=20, color=C_MID, align=PP_ALIGN.CENTER)

# タイトル区切り線
add_rect(p1, Inches(1.5), Inches(3.6), Inches(5.27), Pt(1.5), fill=C_ACCENT)

add_text(p1, "修繕委員会", Inches(0.6), Inches(3.7), Inches(7.07), Inches(0.75),
         size=32, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
add_text(p1, "第1回定例（キックオフ）", Inches(0.6), Inches(4.4), Inches(7.07), Inches(0.65),
         size=24, bold=False, color=C_MID, align=PP_ALIGN.CENTER)

add_rect(p1, Inches(1.5), Inches(5.1), Inches(5.27), Pt(1), fill=C_LIGHT)

add_text(p1, "2026年6月7日（日）　14:00〜15:00",
         Inches(0.6), Inches(5.2), Inches(7.07), Inches(0.5),
         size=16, color=C_SUB, align=PP_ALIGN.CENTER)
add_text(p1, "三鷹市下連雀八丁目地区公会堂",
         Inches(0.6), Inches(5.68), Inches(7.07), Inches(0.45),
         size=14, color=C_SUB, align=PP_ALIGN.CENTER)

# 出席メンバー
add_rect(p1, Inches(1.2), Inches(6.5), Inches(5.87), Inches(0.65),
         fill=C_WHITE, line_color=C_LIGHT, line_pt=1.0)
add_text(p1, "出席：小橋　／　石川　／　十文字　／　阿部　／　武藤　／　新関",
         Inches(1.2), Inches(6.6), Inches(5.87), Inches(0.5),
         size=14, color=C_DARK, align=PP_ALIGN.CENTER)

# 下部バー
add_rect(p1, 0, Inches(11.51), W, Inches(0.18), fill=C_ACCENT)

# ══════════════════════════════════════════════════════
# PAGE 2 ｜ 委員会の目的・アジェンダ
# ══════════════════════════════════════════════════════
p2 = prs.slides.add_slide(BLANK)
add_rect(p2, 0, 0, W, Inches(0.18), fill=C_ACCENT)
add_text(p2, "委員会の目的・本日のアジェンダ",
         Inches(0.5), Inches(0.25), Inches(7.27), Inches(0.65),
         size=22, bold=True, color=C_DARK)
add_rect(p2, Inches(0.5), Inches(0.9), Inches(7.27), Pt(1.5), fill=C_ACCENT)
page_footer(p2, 2)

# 目的
y = section_title(p2, "委員会の目的", Inches(1.05))
items = [
    "● 建物・設備の状態を把握し、長期修繕計画の策定と\n　 必要に応じた修繕積立金の見直しを推進",
    "● 住民の意見を反映した計画立案と情報共有を支援",
    "● 決定は理事会・総会が行い、委員会は調査・提案・情報発信を担う",
]
y = bullet(p2, items, Inches(0.6), y + Inches(0.1), Inches(7.1), size=14, gap=Inches(0.7))

# 位置づけ
add_rect(p2, Inches(0.5), y + Inches(0.1), Inches(7.27), Inches(0.52),
         fill=C_PALE, line_color=C_LIGHT, line_pt=0.75)
add_text(p2, "位置づけ：理事会と連携する住民委員会　／　決定機関は理事会・総会",
         Inches(0.65), y + Inches(0.18), Inches(7.0), Inches(0.38),
         size=13, color=C_SUB, italic=True)
y += Inches(0.82)

# アジェンダ
y = section_title(p2, "本日のアジェンダ", y + Inches(0.2))
agenda = [
    ("1", "メンバー紹介・役割分担"),
    ("2", "活動スコープ・スケジュールの確認"),
    ("3", "コンサル業者選定の進め方"),
    ("4", "省エネ改修・緊急修繕の総会提案（新テーマ）"),
    ("5", "運営ルール・決定事項・次回予定"),
]
y += Inches(0.1)
for num, item in agenda:
    bg = C_PALE if int(num) % 2 == 1 else C_WHITE
    add_rect(p2, Inches(0.5), y, Inches(7.27), Inches(0.52),
             fill=bg, line_color=C_LIGHT, line_pt=0.5)
    add_rect(p2, Inches(0.5), y, Inches(0.5), Inches(0.52), fill=C_ACCENT)
    add_text(p2, num, Inches(0.5), y + Inches(0.08), Inches(0.5), Inches(0.4),
             size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(p2, item, Inches(1.1), y + Inches(0.1), Inches(6.6), Inches(0.38),
             size=15, color=C_DARK)
    y += Inches(0.56)

# ══════════════════════════════════════════════════════
# PAGE 3 ｜ メンバー紹介・役割分担
# ══════════════════════════════════════════════════════
p3 = prs.slides.add_slide(BLANK)
add_rect(p3, 0, 0, W, Inches(0.18), fill=C_ACCENT)
add_text(p3, "メンバー紹介・役割分担",
         Inches(0.5), Inches(0.25), Inches(7.27), Inches(0.65),
         size=22, bold=True, color=C_DARK)
add_rect(p3, Inches(0.5), Inches(0.9), Inches(7.27), Pt(1.5), fill=C_ACCENT)
page_footer(p3, 3)

# メンバー
y = section_title(p3, "出席メンバー", Inches(1.05))
add_rect(p3, Inches(0.5), y + Inches(0.1), Inches(7.27), Inches(0.65),
         fill=C_PALE, line_color=C_LIGHT, line_pt=0.75)
add_text(p3, "小橋　　石川　　十文字　　阿部　　武藤　　新関",
         Inches(0.5), y + Inches(0.15), Inches(7.27), Inches(0.55),
         size=18, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
y += Inches(1.0)

# 役割分担
y = section_title(p3, "【議題】役割分担の決定", y + Inches(0.25))
y += Inches(0.1)
add_table(p3,
    ["役割", "内容"],
    [
        ["委員長",   "委員会の取りまとめ・対外窓口・理事会への報告"],
        ["副委員長", "委員長補佐・議事録管理（AIツール活用）・情報共有の整理"],
    ],
    Inches(0.5), y, Inches(7.27), Inches(1.35),
    hdr_size=14, data_size=14,
    col_widths=[Inches(1.6), Inches(5.67)]
)
y += Inches(1.55)

# 追加担当議論
note_box(p3, "追加担当について（本日議論）\n"
             "● 積立金評価担当（修繕積立金の適正性検証）を設けるか\n"
             "● その他、必要な担当・役割があるかを議論する",
         Inches(0.5), y, Inches(7.27), Inches(1.3))

# ══════════════════════════════════════════════════════
# PAGE 4 ｜ 活動スコープ・改定スケジュール
# ══════════════════════════════════════════════════════
p4 = prs.slides.add_slide(BLANK)
add_rect(p4, 0, 0, W, Inches(0.18), fill=C_ACCENT)
add_text(p4, "活動スコープ・改定スケジュール",
         Inches(0.5), Inches(0.25), Inches(7.27), Inches(0.65),
         size=22, bold=True, color=C_DARK)
add_rect(p4, Inches(0.5), Inches(0.9), Inches(7.27), Pt(1.5), fill=C_ACCENT)
page_footer(p4, 4)

# スコープ
y = section_title(p4, "活動スコープ（3本柱）", Inches(1.05))
scope = [
    ("①", "長期修繕計画の策定（30年）",
     "コンサル選定 → 現地調査 → 計画書作成 → 総会議決"),
    ("②", "短期修繕計画の立案（直近5年）",
     "優先修繕箇所の整理 → コンサル提案 → 総会議決"),
    ("③", "省エネ改修・緊急修繕の提案【NEW】",
     "断熱窓を起点に → 積立金との整合性を確認しながら来年度総会への提案を並行推進"),
]
y += Inches(0.1)
for num, title, detail in scope:
    bg = C_AMBER if num == "③" else C_PALE
    add_rect(p4, Inches(0.5), y, Inches(7.27), Inches(0.68),
             fill=bg, line_color=C_LIGHT, line_pt=0.5)
    add_rect(p4, Inches(0.5), y, Inches(0.46), Inches(0.68), fill=C_ACCENT)
    add_text(p4, num, Inches(0.5), y+Inches(0.1), Inches(0.46), Inches(0.5),
             size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(p4, f"{title}", Inches(1.08), y+Inches(0.04), Inches(6.6), Inches(0.3),
             size=14, bold=True, color=C_DARK)
    add_text(p4, detail, Inches(1.08), y+Inches(0.34), Inches(6.6), Inches(0.32),
             size=12, color=C_SUB)
    y += Inches(0.74)

# スケジュール注記
add_text(p4,
    "※キックオフが5月予定から6月7日に遅延。以下は1ヶ月後ろ倒しの改定版。",
    Inches(0.5), y + Inches(0.2), Inches(7.27), Inches(0.36),
    size=12, color=C_WARN, italic=True)
y += Inches(0.72)

# ①長期スケジュール
y = section_title(p4, "① 長期修繕計画スケジュール", y)
y += Inches(0.08)
long_rows = [
    ["6月〜",  "コンサル3社へプレゼン依頼・日程調整"],
    ["7月",    "コンサル3社プレゼン実施（TDS・オフィスレコン・センターオフィス）"],
    ["8月",    "業者選定・理事会決議"],
    ["9月",    "実施計画案の提示"],
    ["10月",   "実施計画案の決議・現地調査（委員会立会い）"],
    ["11月",   "調査報告・中長期修繕計画 納品・理事会報告"],
    ["12月",   "住民向け全戸配布 / 修繕積立金変更案の検討"],
    ["来年総会","修繕計画・積立金変更の議決"],
]
add_table(p4, ["時期", "内容（委員会中心）"],
          long_rows, Inches(0.5), y, Inches(7.27), Inches(3.3),
          hdr_size=13, data_size=13,
          col_widths=[Inches(1.15), Inches(6.12)])

# ══════════════════════════════════════════════════════
# PAGE 5 ｜ 短期スケジュール・省エネ改修
# ══════════════════════════════════════════════════════
p5 = prs.slides.add_slide(BLANK)
add_rect(p5, 0, 0, W, Inches(0.18), fill=C_ACCENT)
add_text(p5, "活動スコープ・改定スケジュール（続き）",
         Inches(0.5), Inches(0.25), Inches(7.27), Inches(0.65),
         size=22, bold=True, color=C_DARK)
add_rect(p5, Inches(0.5), Inches(0.9), Inches(7.27), Pt(1.5), fill=C_ACCENT)
page_footer(p5, 5)

# ②短期スケジュール
y = section_title(p5, "② 短期修繕計画スケジュール（直近5年）", Inches(1.05))
y += Inches(0.08)
short_rows = [
    ["8月",    "要求仕様提示（インフラ緊急性・省エネ優先）"],
    ["9月",    "計画案提示（省エネ改修を中心に）"],
    ["11月",   "現地調査後の優先修繕案提示"],
    ["12月",   "短期修繕計画案の検討 / 総会提案資料作成"],
    ["来年総会","議決"],
]
add_table(p5, ["時期", "内容（委員会中心）"],
          short_rows, Inches(0.5), y, Inches(7.27), Inches(2.15),
          hdr_size=13, data_size=13,
          col_widths=[Inches(1.15), Inches(6.12)])
y += Inches(2.35)

# 省エネ改修・緊急修繕
y = section_title(p5, "新テーマ：省エネ改修・緊急修繕の来年総会提案", y + Inches(0.2))

add_rect(p5, Inches(0.5), y + Inches(0.05), Inches(7.27), Inches(0.45),
         fill=C_PALE, line_color=C_LIGHT, line_pt=0.5)
add_text(p5,
    "長期計画策定と並行し、早期に独立した検討ラインで動く。各提案は修繕積立金への影響を確認しながら進める。",
    Inches(0.65), y + Inches(0.1), Inches(7.0), Inches(0.38),
    size=13, color=C_SUB, italic=True)
y += Inches(0.65)

add_table(p5,
    ["テーマ", "内容", "目標"],
    [
        ["省エネ\n改修",
         "断熱窓改修を優先検討（専有部対象）\n補助金・助成金（都・市・国）の活用可能性を調査\nコンサルにも提案を依頼",
         "2027年\n総会への\n提案"],
        ["緊急\n修繕",
         "インフラ系（給排水管・電気設備等）の現状把握を優先\nコンサル調査と並行して早期実態把握\n緊急度が高いものは総会を待たず理事会決議で対応",
         "随時\n理事会\n判断"],
    ],
    Inches(0.5), y, Inches(7.27), Inches(2.6),
    hdr_size=13, data_size=13,
    col_widths=[Inches(0.9), Inches(5.47), Inches(0.9)]
)
y += Inches(2.8)

add_rect(p5, Inches(0.5), y, Inches(7.27), Inches(0.6),
         fill=C_PALE, line_color=C_LIGHT, line_pt=0.5)
add_text(p5,
    "● 省エネ改修：補助金差引後の実質負担と積立金残高を照合\n"
    "● 緊急修繕：修繕費用の積立金への影響を試算・確認（積立金評価担当設置の場合）",
    Inches(0.65), y + Inches(0.06), Inches(7.0), Inches(0.52),
    size=13, color=C_SUB)

# ══════════════════════════════════════════════════════
# PAGE 6 ｜ コンサル業者選定の進め方
# ══════════════════════════════════════════════════════
p6 = prs.slides.add_slide(BLANK)
add_rect(p6, 0, 0, W, Inches(0.18), fill=C_ACCENT)
add_text(p6, "コンサル業者選定の進め方",
         Inches(0.5), Inches(0.25), Inches(7.27), Inches(0.65),
         size=22, bold=True, color=C_DARK)
add_rect(p6, Inches(0.5), Inches(0.9), Inches(7.27), Pt(1.5), fill=C_ACCENT)
page_footer(p6, 6)

# 候補3社
add_rect(p6, Inches(0.5), Inches(1.05), Inches(7.27), Inches(0.55),
         fill=C_PALE, line_color=C_ACCENT, line_pt=1.0)
add_text(p6, "候補3社：　TDS　／　オフィスレコン　／　センターオフィス",
         Inches(0.65), Inches(1.12), Inches(7.0), Inches(0.45),
         size=16, bold=True, color=C_MID, align=PP_ALIGN.CENTER)

y = Inches(1.75)
y = section_title(p6, "プレゼン実施方針", y)
items = [
    "● 7月中に3社それぞれ30分のプレゼンを実施",
    "● 場所：会場を1日確保して3社連続、または管理事務室（Web参加可）",
    "● 本日：各社への連絡・日程調整・会場確保の担当を決定",
]
y = bullet(p6, items, Inches(0.6), y + Inches(0.08), Inches(7.1), size=14, gap=Inches(0.5))

y = section_title(p6, "選定基準（3点）", y + Inches(0.15))
y += Inches(0.08)
add_table(p6,
    ["#", "基準"],
    [
        ["①", "要求仕様を満たしているか"],
        ["②", "金額が妥当か（想定100〜200万円）"],
        ["③", "短期計画・省エネ改修のコンサル提案が可能か"],
    ],
    Inches(0.5), y, Inches(7.27), Inches(1.5),
    hdr_size=14, data_size=14,
    col_widths=[Inches(0.45), Inches(6.82)]
)
y += Inches(1.7)

y = section_title(p6, "依頼内容", y)
items2 = [
    "● 劣化診断調査（第3回大規模修繕に向けた建物調査）",
    "● 30年中長期修繕計画の作成・修繕積立金の評価",
    "● 短期修繕計画（5年）の立案（追加依頼）",
    "● 省エネ改修提案（追加依頼）",
]
y = bullet(p6, items2, Inches(0.6), y + Inches(0.08), Inches(7.1), size=14, gap=Inches(0.5))

# 議題ボックス
note_box(p6,
    "【議題】見積事業者選定プロセスの追加について\n"
    "● 管理会社任せにせず、委員会として選定基準・比較評価プロセスを持つことで透明性が高まる\n"
    "● 一方、現時点では計画策定が優先。コンサル選定後に整備する選択肢もある\n"
    "● 論点：今の段階で仕組みを検討するか、後回しにするか",
    Inches(0.5), y + Inches(0.2), Inches(7.27), Inches(1.55))

# ══════════════════════════════════════════════════════
# PAGE 7 ｜ 運営ルール・決定事項・次回予定
# ══════════════════════════════════════════════════════
p7 = prs.slides.add_slide(BLANK)
add_rect(p7, 0, 0, W, Inches(0.18), fill=C_ACCENT)
add_text(p7, "運営ルール・本日の決定事項・次回予定",
         Inches(0.5), Inches(0.25), Inches(7.27), Inches(0.65),
         size=22, bold=True, color=C_DARK)
add_rect(p7, Inches(0.5), Inches(0.9), Inches(7.27), Pt(1.5), fill=C_ACCENT)
page_footer(p7, 7)

y = section_title(p7, "運営ルール", Inches(1.05))
y += Inches(0.08)
add_table(p7,
    ["項目", "内容"],
    [
        ["打合せ頻度", "月1回程度"],
        ["コミュニケーション", "Google Meet・チャット＋メール、必要に応じ書面・対面"],
        ["議事録", "副委員長（AIツール活用）"],
    ],
    Inches(0.5), y, Inches(7.27), Inches(1.6),
    hdr_size=14, data_size=14,
    col_widths=[Inches(2.0), Inches(5.27)]
)
y += Inches(1.8)

y = section_title(p7, "本日の決定事項", y)
decisions = [
    "□　役割分担（委員長・副委員長・その他担当）",
    "□　コンサル3社への連絡・プレゼン日程調整・会場確保の担当",
    "□　見積事業者選定プロセスを今設けるか後回しにするか",
    "□　省エネ改修・緊急修繕の優先確認事項",
    "□　次回開催日程",
]
y += Inches(0.08)
for i, d in enumerate(decisions):
    bg = C_PALE if i % 2 == 0 else C_WHITE
    add_rect(p7, Inches(0.5), y, Inches(7.27), Inches(0.5),
             fill=bg, line_color=C_LIGHT, line_pt=0.5)
    add_text(p7, d, Inches(0.68), y + Inches(0.08), Inches(7.0), Inches(0.38),
             size=14, color=C_DARK)
    y += Inches(0.5)

y += Inches(0.2)
y = section_title(p7, "次回（7月）　コンサル3社プレゼン実施　→　業者選定の議論", y)
items3 = [
    "● コンサル3社（TDS・オフィスレコン・センターオフィス）のプレゼンを実施（各30分）",
    "● 各社の提案内容・金額・対応範囲を委員会内で評価・比較",
    "● 業者選定の方向性を議論し、8月の理事会決議へ向けて準備",
]
bullet(p7, items3, Inches(0.6), y + Inches(0.08), Inches(7.1), size=14, gap=Inches(0.52))

# ════════════════════════════════════════════════
out = "/home/user/Ryuichi_Ishikawa_Private/notes/mansion/修繕委員会_第1回定例_キックオフ.pptx"
prs.save(out)
print("saved:", out)
